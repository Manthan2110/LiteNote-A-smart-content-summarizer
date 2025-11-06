import streamlit as st
import google.generativeai as genai
from youtube_transcript_api import YouTubeTranscriptApi
import trafilatura
import requests
from newspaper import Article
from bs4 import BeautifulSoup
from langdetect import detect
from fpdf import FPDF
from docx import Document
from pptx import Presentation
from streamlit_extras.colored_header import colored_header
from streamlit_extras.add_vertical_space import add_vertical_space
import time
from urllib.parse import urlparse
import re
import tempfile


# Page Config
st.set_page_config(
    page_title="LiteNote AI Summarizer",
    page_icon="",
    layout="wide"
)


# Custom CSS - Using the modern design from the second project
st.markdown("""
<style>
body {
    background-color: #f8fafc; 
    font-family: 'Inter', 'Segoe UI', sans-serif;
}
.title {
    text-align: center; 
    font-size: 2.5rem; 
    font-weight: bold; 
    background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
    -webkit-background-clip: text; 
    -webkit-text-fill-color: transparent;
    margin-bottom: 1rem;
}
.subtitle {
    text-align: center;
    color: #64748b;
    font-size: 1.1rem;
    margin-bottom: 2rem;
}
.content-box {
    padding: 1.5rem; 
    border-radius: 12px; 
    box-shadow: 0 4px 6px -1px rgba(0, 0, 0, 0.1);
    font-size: 1rem; 
    line-height: 1.7;
    border: 1px solid #e2e8f0;
}
@media (prefers-color-scheme: light) {
    .content-box {
        background: #ffffff; 
        color: #1a1a1a;
    }
}
@media (prefers-color-scheme: dark) {
    .content-box {
        background: #1e293b; 
        color: #f1f5f9;
    }
}
.stButton>button {
    border-radius: 8px; 
    background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); 
    color: white; 
    font-weight: 600; 
    padding: 0.75rem 1.5rem;
    border: none;
    transition: all 0.3s ease;
}
.stButton>button:hover {
    background: linear-gradient(135deg, #5a67d8 0%, #6b46c1 100%); 
    transform: translateY(-2px);
    box-shadow: 0 4px 12px rgba(102, 126, 234, 0.4);
}
.extraction-method {
    background: #f0f9ff;
    border: 1px solid #0ea5e9;
    border-radius: 6px;
    padding: 0.5rem;
    font-size: 0.9rem;
    color: #0369a1;
    margin: 0.5rem 0;
}
.api-key-info {
    background: #fef3c7;
    border: 1px solid #f59e0b;
    border-radius: 6px;
    padding: 0.75rem;
    font-size: 0.9rem;
    color: #92400e;
    margin: 1rem 0;
}
</style>
""", unsafe_allow_html=True)


# Session state initialization
if 'api_key_set' not in st.session_state:
    st.session_state.api_key_set = False


def is_valid_url(url):
    """Validate URL format"""
    try:
        result = urlparse(url)
        return all([result.scheme, result.netloc])
    except:
        return False


def is_youtube_url(url):
    """Check if URL is a YouTube video"""
    youtube_domains = ['youtube.com', 'youtu.be', 'www.youtube.com', 'm.youtube.com']
    try:
        parsed = urlparse(url)
        return any(domain in parsed.netloc for domain in youtube_domains) and ('v=' in url or 'youtu.be/' in url)
    except:
        return False


def clean_text(text):
    """Clean extracted text"""
    if not text:
        return ""
    text = re.sub(r'\n+', '\n', text)
    text = re.sub(r' +', ' ', text)
    return text.strip()


language_fallbacks = [
    'en', 'hi', 'es', 'fr', 'de', 'ja', 'as', 'bn', 'gu', 'kn', 'ml', 'mr', 'or', 'pa', 'ta', 'te', 'ur'
]


@st.cache_data(show_spinner=False, ttl=3600)
def extract_youtube_transcript(youtube_video_url):
    try:
        if 'youtu.be/' in youtube_video_url:
            video_id = youtube_video_url.split('youtu.be/')[1].split('?')[0]
        else:
            video_id = youtube_video_url.split("v=")[1].split("&")[0]

        yt_api = YouTubeTranscriptApi()
        transcript_list = yt_api.fetch(video_id, languages=language_fallbacks)
        transcript = " ".join([entry.text for entry in transcript_list])

        return {
            'content': clean_text(transcript),
            'video_id': video_id,
            'title': f"YouTube Video ({video_id})",
            'method': 'YouTube Transcript API'
        }
    except Exception as e:
        raise Exception(f"Failed to extract YouTube transcript: {str(e)}")


@st.cache_data(show_spinner=False, ttl=3600)
def extract_content_trafilatura(url):
    try:
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }

        session = requests.Session()
        session.headers.update(headers)
        response = session.get(url, timeout=10)
        html_content = response.text

        extracted = trafilatura.extract(
            html_content,
            include_comments=False,
            include_tables=True,
            include_images=False,
            with_metadata=True,
            output_format='xml'
        )

        if extracted:
            from xml.etree import ElementTree as ET
            try:
                root = ET.fromstring(extracted)
                title = root.find('.//title')
                title = title.text if title is not None else "Unknown Title"

                author = root.find('.//author')
                author = author.text if author is not None else "Unknown Author"

                date = root.find('.//date')
                date = date.text if date is not None else "Unknown Date"

                main_text = trafilatura.extract(html_content, include_comments=False)

                return {
                    'content': clean_text(main_text),
                    'title': title,
                    'author': author,
                    'date': date,
                    'method': 'Trafilatura'
                }
            except:
                main_text = trafilatura.extract(html_content, include_comments=False)
                return {
                    'content': clean_text(main_text),
                    'title': 'Unknown Title',
                    'author': 'Unknown Author',
                    'date': 'Unknown Date',
                    'method': 'Trafilatura'
                }
    except Exception:
        return None


@st.cache_data(show_spinner=False, ttl=3600)
def extract_content_newspaper(url):
    try:
        article = Article(url)
        article.download()
        article.parse()

        if article.text:
            return {
                'content': clean_text(article.text),
                'title': article.title,
                'author': ', '.join(article.authors),
                'date': str(article.publish_date),
                'method': 'Newspaper3k'
            }
    except Exception:
        return None


@st.cache_data(show_spinner=False, ttl=3600)
def extract_content_beautifulsoup(url):
    try:
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }

        response = requests.get(url, headers=headers, timeout=10)
        soup = BeautifulSoup(response.content, 'html.parser')

        title = soup.find('title')
        title = title.get_text().strip() if title else "Unknown Title"

        content_selectors = [
            'article', 'main', '[role="main"]',
            '.content', '.post-content', '.entry-content',
            '.article-body', '.post-body'
        ]

        content = ""
        for selector in content_selectors:
            elements = soup.select(selector)
            if elements:
                content = ' '.join([elem.get_text() for elem in elements])
                break

        if not content:
            paragraphs = soup.find_all('p')
            content = ' '.join([p.get_text() for p in paragraphs])

        if content:
            return {
                'content': clean_text(content),
                'title': title,
                'author': "Unknown Author",
                'date': "Unknown Date",
                'method': 'Beautiful Soup'
            }
    except Exception:
        return None


def extract_website_content(url):
    methods = [
        extract_content_trafilatura,
        extract_content_newspaper,
        extract_content_beautifulsoup
    ]

    for method in methods:
        result = method(url)
        if result and result['content']:
            return result
    return None


def generate_gemini_summary(content_data, lang_choice, summary_level, summary_style, content_type):
    try:
        model = genai.GenerativeModel("gemini-1.5-flash")

        content = content_data['content']
        title = content_data.get('title', 'Unknown Title')

        try:
            detected_lang = detect(content[:500])
        except:
            detected_lang = "unknown"

        if lang_choice == "Auto (Content Language)":
            lang_instruction = f"The content is in **{detected_lang}**. Summarize in the same language."
        else:
            lang_instruction = f"Translate and summarize into **{lang_choice}**."

        length_instructions = {
            "Brief": "Summarize concisely in 3-5 bullet points per section.",
            "Medium": "Summarize in 6-10 bullet points or short paragraphs per section.",
            "Detailed": "Provide detailed paragraphs with comprehensive explanations, examples, and context for each section."
        }

        style_instruction = "Use bullet points." if summary_style == "Bullets" else "Use paragraph format."

        if content_type == "youtube":
            type_instruction = """
            You are summarizing a YouTube video transcript. Focus on:
            - Main topics and key messages from the video
            - Important tips, insights, or tutorials mentioned
            - Sequential flow of information as presented in the video
            """
            source_info = f"""
            **Source Information:**
            - Video Title: {title}
            - Content Type: YouTube Video Transcript
            - Original Language: {detected_lang}
            """
        else:
            type_instruction = """
            You are summarizing web content from an article or blog. Focus on:
            - Main arguments and key points
            - Supporting evidence and data
            - Conclusions and recommendations
            """
            author = content_data.get('author', 'Unknown Author')
            source_info = f"""
            **Source Information:**
            - Title: {title}
            - Author: {author}
            - Content Type: Website/Blog Article
            - Original Language: {detected_lang}
            """

        prompt = f"""
        You are a highly skilled multilingual content summarizer and analyst.

        {lang_instruction}
        {length_instructions[summary_level]}
        {style_instruction}
        {type_instruction}

        Your task is to create a **structured, comprehensive, and actionable summary** of the provided content.

        {source_info}

        **Instructions:**

        1. **Document Header**
           - Include the source information above
           - Summary Language: {lang_choice}

        2. **Executive Summary**
           - Provide a 2-3 sentence overview of the main topic and key findings

        3. **Main Content Analysis**
           - Identify and organize content into logical sections
           - Use descriptive headings for each section
           - Extract key insights, arguments, and supporting evidence
           - Highlight important data, statistics, or quotes (if present)

        4. **Key Takeaways Table**
           Create a markdown table:
           | Section | Key Insight |
           |---------|-------------|
           | Section Name | One-sentence takeaway |

        5. **Actionable Insights**
           - List 3-5 practical insights or recommendations
           - Make them specific and directly applicable

        6. **Content Assessment**
           - Brief note on content quality, credibility, and usefulness

        **Formatting Requirements:**
        - Use clear Markdown formatting
        - Maintain logical flow and readability
        - Avoid repetition and filler content
        - Focus on value-driven insights

        **Content to Summarize:**
        {content[:8000]}
        """

        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        st.error(f"Summary generation failed: {str(e)}")
        return None


def create_download_files(summary, title):
    try:
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=11)
        summary_text = summary.encode('latin-1', 'replace').decode('latin-1')
        pdf.multi_cell(0, 8, summary_text)
        pdf_output = pdf.output(dest='S').encode('latin-1')

        doc = Document()
        doc.add_heading('Content Summary', 0)
        doc.add_heading(f'Title: {title}', level=2)
        doc.add_paragraph(summary)

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Content Summary"
        slide.placeholders[1].text = summary[:500] + "..." if len(summary) > 500 else summary

        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp_docx:
            doc.save(tmp_docx.name)
            with open(tmp_docx.name, 'rb') as f:
                docx_bytes = f.read()

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_pptx:
            prs.save(tmp_pptx.name)
            with open(tmp_pptx.name, 'rb') as f:
                pptx_bytes = f.read()

        return pdf_output, docx_bytes, pptx_bytes
    except Exception as e:
        st.warning(f"Some download formats may not be available: {str(e)}")
        return None, None, None


# Sidebar - API Key input and info
with st.sidebar:
    st.markdown("### 🔐 Setup Google Gemini API")
    api_key = st.text_input(
        "Enter your Google Gemini API Key:",
        type="password",
        help="Get your API key from Google AI Studio: https://makersuite.google.com/app/apikey"
    )
    if st.button("🚀 Set API Key"):
        if api_key:
            try:
                genai.configure(api_key=api_key)
                model = genai.GenerativeModel("gemini-1.5-flash")
                test_response = model.generate_content("Hello")
                st.session_state.api_key_set = True
                st.success("✅ API Key configured successfully!")
                st.experimental_rerun()
            except Exception as e:
                st.error(f"❌ Invalid API Key: {str(e)}")
        else:
            st.warning("⚠️ Please enter your API key")
    st.markdown("""
    <div class="api-key-info">
    🔒 Your API key is used only for this session. Get one at 
    <a href="https://makersuite.google.com/app/apikey" target="_blank">Google AI Studio</a>.
    </div>
    """, unsafe_allow_html=True)
    add_vertical_space(1)
    st.markdown("### ℹ️ About")
    st.markdown("Powered by **Google Gemini AI** with advanced content extraction capabilities.")


# Main interface header
st.markdown("<h1 class='title'>LiteNote - AI Content Summarizer</h1>", unsafe_allow_html=True)
st.markdown("<p class='subtitle'>Transform YouTube videos and web content into structured, actionable insights</p>", unsafe_allow_html=True)


if not st.session_state.api_key_set:
    st.warning("🔐 Please enter and set your Google Gemini API key in the sidebar to use this app.")
    st.stop()


# Single input box for both YouTube and Website/Blog URLs
input_url = st.text_input("Enter YouTube video or Website/Blog URL:", placeholder="Paste your link here...")

if input_url:
    if not is_valid_url(input_url):
        st.error("❌ Please enter a valid URL starting with http:// or https://")
        st.stop()

    col1, col2 = st.columns([1,1])
    with col1:
        lang_choice = st.selectbox("Summary Language",
                                  ["Auto (Content Language)", "English", "Hindi", "Spanish", "French", "German"])
    with col2:
        summary_level = st.selectbox("Summary Length", ["Brief", "Medium", "Detailed"])
    summary_style = st.radio("Summary Style", ["Bullets", "Paragraphs"], horizontal=True)

    if st.button("✨ Generate Summary"):
        with st.spinner("Extracting content and generating summary..."):
            try:
                if is_youtube_url(input_url):
                    extracted_data = extract_youtube_transcript(input_url)
                    content_type = "youtube"
                else:
                    extracted_data = extract_website_content(input_url)
                    content_type = "website"

                if not extracted_data or not extracted_data.get('content'):
                    st.error("❌ Failed to extract content. Please check the URL and try again.")
                    st.stop()

                summary = generate_gemini_summary(
                    extracted_data, lang_choice, summary_level, summary_style, content_type
                )

                if not summary:
                    st.error("❌ Failed to generate summary. Please try again.")
                    st.stop()

                colored_header("📄 Summary", color_name="blue-70")
                st.markdown(f"<div class='extraction-method'>✅ Content extracted using: <strong>{extracted_data.get('method', 'Unknown')}</strong></div>", unsafe_allow_html=True)
                st.markdown(f"<div class='content-box'>{summary}</div>", unsafe_allow_html=True)

                # Download options
                st.markdown("### 📥 Download Summary")
                pdf_data, docx_data, pptx_data = create_download_files(summary, extracted_data.get('title', 'Summary'))

                col1, col2, col3, col4 = st.columns(4)
                with col1:
                    st.download_button("📄 TXT", summary, "summary.txt", "text/plain")
                if pdf_data:
                    with col2:
                        st.download_button("📕 PDF", pdf_data, "summary.pdf", "application/pdf")
                if docx_data:
                    with col3:
                        st.download_button("📘 Word", docx_data, "summary.docx",
                                            "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
                if pptx_data:
                    with col4:
                        st.download_button("📊 PowerPoint", pptx_data, "summary.pptx",
                                            "application/vnd.openxmlformats-officedocument.presentationml.presentation")

            except Exception as e:
                st.error(f"❌ Error during processing: {str(e)}")

# Footer
st.markdown("---")
st.markdown("""
<div style='text-align: center; color: #64748b;'>
    <p>🤖 Powered by Google Gemini AI | 🛠️ Built with Streamlit</p>
    <p><small>Universal AI-powered content summarization for YouTube videos and web articles</small></p>
</div>
""", unsafe_allow_html=True)

