import streamlit as st
from dotenv import load_dotenv
import os
import google.generativeai as genai
from youtube_transcript_api import YouTubeTranscriptApi
from langdetect import detect
from fpdf import FPDF
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from streamlit_extras.colored_header import colored_header
from streamlit_extras.add_vertical_space import add_vertical_space

# Load environment variables
load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

# Page Config
st.set_page_config(page_title="YT Notes Generator", page_icon="🎬", layout="wide")

# Custom CSS
st.markdown("""
<style>
body {background-color: #f9fafc; font-family: 'Segoe UI', sans-serif;}
.title {text-align: center; font-size: 2.2rem; font-weight: bold; background: linear-gradient(90deg, #ff6a00, #ee0979);
-webkit-background-clip: text; -webkit-text-fill-color: transparent;}
.summary-box {padding:1.2rem; border-radius:15px; box-shadow:0 4px 20px rgba(0,0,0,0.08); font-size:1rem; line-height:1.6;}
@media (prefers-color-scheme: light) {.summary-box {background: #ffffff; color: #1a1a1a;}}
@media (prefers-color-scheme: dark) {.summary-box {background: #1e1e1e; color: #f5f5f5;}}

.stButton>button {border-radius:10px; background:linear-gradient(90deg,#667eea,#764ba2); color:white; font-weight:bold; padding:0.6rem 1.2rem;}
.stButton>button:hover {background:linear-gradient(90deg,#764ba2,#667eea); color:white;}
</style>
""", unsafe_allow_html=True)

# Indian + global language fallback list
language_fallbacks = [
    'en', 'hi', 'es', 'fr', 'de', 'ja', 'as', 'bn', 'gu', 'kn', 'ml', 'mr', 'or', 'pa', 'ta', 'te', 'ur'
]

# Cache transcript
@st.cache_data(show_spinner=False)
def extract_transcript_details(youtube_video_url):
    try:
        video_id = youtube_video_url.split("v=")[1].split("&")[0]
        yt_api = YouTubeTranscriptApi()
        transcript_list = yt_api.fetch(video_id, languages=language_fallbacks)
        transcript = " ".join([entry.text for entry in transcript_list])
        return transcript, video_id
    except Exception as e:
        raise e

# Generate Gemini content
def generate_gemini_content(transcript_text, lang_choice, summary_level, summary_style):
    model = genai.GenerativeModel("gemini-2.5-flash")
    detected_lang = detect(transcript_text[:500])

    # Language instruction
    if lang_choice == "Auto (Transcript Language)":
        lang_instruction = f"The transcript is in **{detected_lang}**. Summarize in the same language."
    else:
        lang_instruction = f"Translate and summarize into **{lang_choice}**."

    # Summary length & style instruction
    if summary_level == "Brief":
        length_instruction = "Summarize concisely in 3–5 bullet points per section."
    elif summary_level == "Medium":
        length_instruction = "Summarize in 6–10 bullet points or short paragraphs per section."
    else:
        length_instruction = "Provide detailed paragraphs with comprehensive explanations, examples, and context for each section."

    style_instruction = "Use bullet points." if summary_style == "Bullets" else "Use paragraph format."

    prompt = f"""
    You are a highly skilled multilingual YouTube transcript summarizer.
    {lang_instruction}
    {length_instruction}
    {style_instruction}

    Your goal is to generate a **structured, readable, and actionable summary**. Follow these detailed instructions:

    1. **Title & Metadata**  
    - Include the video title and channel name both in different line (if available).  
    - Mention the original language of the transcript.
    - Like this example:
        Video Title: The Future of AI  
        Channel: Tech Insights  
        Original Language: English

    2. **Sections & Headings**  
    - Identify and divide content into key sections or topics.  
    - Use descriptive headings for each section.  
    - Ensure logical flow: introduction → main content → conclusion.

    3. **Section Summaries**  
    - For each section, provide concise, clear bullet points or paragraphs.  
    - Highlight important tips, examples, explanations, or quotes (only if present in transcript).  
    - Avoid repetition and filler words.

    4. **Summary Table**  
    - At the end, create a Markdown table listing:  
        | Section | Key Takeaway |  
        - One or two sentence takeaway for each section.

    5. **Key Insights & Actionable Points**  
    - Extract 3–5 high-value insights or actionable points from the video.  
    - Make them practical and directly usable by the reader.

    6. **Formatting Guidelines**  
    - Use Markdown for headings, bold/italic text, and lists.  
    - For bullet points: use `-` or `*` consistently.  
    - For paragraphs: maintain clear spacing and readability.

    Output Format: Markdown only.
    """
    response = model.generate_content(prompt + transcript_text)
    return response.text

# Sidebar
with st.sidebar:
    st.image("https://www.gstatic.com/youtube/img/branding/youtubelogo/svg/youtubelogo.svg", width=150)
    st.markdown("### 📌 How it works")
    st.markdown("""
    1. Paste a YouTube link 🎥  
    2. Select output language 🌍  
    3. Select summary length & style  
    4. Click **Get Notes** ✨  
    5. Download summary 📖  
    """)
    add_vertical_space(1)
    st.markdown("### ℹ️ About")
    st.markdown("This app uses **Google Gemini** to generate multilingual structured notes from YouTube transcripts.")

# Title
st.markdown("<h1 class='title'>YouTube Transcript → Smart Notes</h1>", unsafe_allow_html=True)

# Input
youtube_link = st.text_input("🔗 Enter YouTube Video Link:")

# Summary options just after link input
if youtube_link:
    col1, col2 = st.columns([1, 1])
    with col1:
        lang_choice = st.selectbox(
            "Choose summary language:",
            ["Auto (Transcript Language)", "English", "Hindi", "Spanish", "French", "German"]
        )
    with col2:
        summary_level = st.selectbox("Summary Length:", ["Brief", "Medium", "Detailed"])
    summary_style = st.radio("Summary Style:", ["Bullets", "Paragraphs"])

# Show thumbnail
if youtube_link and "v=" in youtube_link:
    video_id = youtube_link.split("v=")[1].split("&")[0]
    st.image(f"http://img.youtube.com/vi/{video_id}/0.jpg", use_column_width=True)

# Generate summary
if st.button("✨ Get Detailed Notes"):
    with st.spinner("⏳ Extracting transcript & generating summary..."):
        try:
            transcript_text, video_id = extract_transcript_details(youtube_link)

            summary = generate_gemini_content(
                transcript_text, lang_choice, summary_level, summary_style
            )

            # Show summary
            colored_header("📑 Detailed Notes", description=None, color_name="violet-70")
            st.markdown(f"<div class='summary-box'>{summary}</div>", unsafe_allow_html=True)

            # Downloads in one column
            st.markdown("### 📥 Download Summary")

            # Create files first
            # PDF
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            pdf.multi_cell(0, 10, summary)
            pdf_file = "YT_Summary.pdf"
            pdf.output(pdf_file)

            # Word
            doc = Document()
            doc.add_heading("YouTube Summary", 0)
            doc.add_paragraph(summary)
            docx_file = "YT_Summary.docx"
            doc.save(docx_file)

            # PPT
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "YouTube Summary"
            slide.placeholders[1].text = summary
            pptx_file = "YT_Summary.pptx"
            prs.save(pptx_file)

            # Display buttons in one line
            col1, col2, col3, col4 = st.columns(4)

            with col1:
                st.download_button("TXT", summary, file_name="YT_Summary.txt")

            with col2:
                st.download_button("PDF", data=open(pdf_file, "rb"), file_name=pdf_file)

            with col3:
                st.download_button("Word", data=open(docx_file, "rb"), file_name=docx_file)

            with col4:
                st.download_button("PPT", data=open(pptx_file, "rb"), file_name=pptx_file)


        except Exception as e:
            st.error(f"❌ Error: {e}")
